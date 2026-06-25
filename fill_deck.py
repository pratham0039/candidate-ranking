"""Fill the mandatory Redrob idea-submission template with our content.

Preserves the template's font (Manrope) and dark colour; only replaces the
prompt text in each content box with our answers. Run:
    python fill_deck.py
Writes 'Idea Submission Template _ Redrob.pptx' in place (back it up first).
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

INK = RGBColor(0x20, 0x27, 0x29)
ACCENT = RGBColor(0x1A, 0x73, 0xE8)
FONT = "Manrope SemiBold"
SRC = "Idea Submission Template _ Redrob.pptx"

REPO = "https://github.com/pratham0039/candidate-ranking"
SANDBOX = "<add Streamlit Cloud link before upload>"


def fill(box, paras):
    """paras: list of (text, size_pt, bold, accent). First fills the
    existing paragraph; the rest are appended."""
    tf = box.text_frame
    tf.word_wrap = True
    # drop all but the first paragraph
    p0 = tf.paragraphs[0]
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

    for i, (text, size, bold, accent) in enumerate(paras):
        para = p0 if i == 0 else tf.add_paragraph()
        para.space_after = Pt(3)
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = ACCENT if accent else INK


def append_value(box, value):
    """Append a value run to a title-slide label paragraph."""
    para = box.text_frame.paragraphs[0]
    src_font = para.runs[0].font if para.runs else None
    run = para.add_run()
    run.text = " " + value
    run.font.name = FONT
    run.font.size = src_font.size if src_font and src_font.size else Pt(14)
    run.font.color.rgb = ACCENT


# (text, size, bold, accent)
H = lambda t: (t, 12.5, True, True)      # section heading
B = lambda t: ("•  " + t, 10.5, False, False)   # bullet
S = lambda t: ("    – " + t, 10, False, False)   # sub-bullet
P = lambda t, s=10.5: (t, s, False, False)        # plain line

CONTENT = {
    # slide index -> shape_id -> list of paragraphs
    1: {64: [
        H("What is our solution?"),
        B("A recruiter-style ranking funnel: parse the JD into a structured "
          "spec, gate out irrelevant and internally-impossible profiles, then "
          "score survivors with a multiplicative blend of factors."),
        B("Hybrid relevance = corroborated career-evidence + TF-IDF + sentence "
          "embeddings. Runs 100,000 to 100 in ~14s, CPU-only, no network."),
        H("What makes it different"),
        B("Scores what candidates DID (career history), not the self-reported "
          "skills list, so keyword stuffing fails."),
        B("Multiplicative scoring: an unreachable or inconsistent candidate "
          "cannot be rescued by a strong skill match alone."),
        B("Every ranked candidate carries an auto-generated, fact-checked "
          "justification."),
    ]},
    2: {71: [
        H("Key requirements extracted from the JD"),
        B("Experience 5-9 yrs (ideal 6-8); Pune/Noida preferred, metros "
          "welcome, no visa sponsorship; notice <=30 days."),
        B("Must-haves: embeddings/retrieval, vector-search infra, ranking "
          "systems, ranking evaluation, strong Python. Nice-to-haves: LLM "
          "fine-tuning, learning-to-rank, HR-tech domain, scale."),
        H("Most important candidate signals"),
        B("Production evidence of retrieval/ranking work in the career "
          "history; experience-band fit; behavioral availability (reply rate, "
          "recency, notice, open-to-work)."),
        H("Evaluating fit beyond keywords"),
        B("A listed skill counts only when corroborated by an assessment "
          "score or real usage duration; embeddings surface strong candidates "
          "who use zero buzzwords."),
    ]},
    3: {78: [
        H("Retrieve, score, rank"),
        B("Funnel: 100,000 -> title gate (-68,821) -> consistency gate "
          "(-482) -> 17,391 scored -> top 100."),
        P("score = skill_fit x experience^w x location^w x availability^w "
          "x ownership^w x penalties^w", 10),
        P("skill_fit = 0.15*evidence + 0.41*TF-IDF + 0.33*embedding", 10),
        H("Models, algorithms, heuristics"),
        B("MiniLM (all-MiniLM-L6-v2) embeddings, scikit-learn TF-IDF, "
          "deterministic feature functions. No LLM at rank time."),
        H("Combining signals into one ranking"),
        B("Multiplicative product of normalized factors; exponents tuned by "
          "random search and minimax over six label hypotheses, against an "
          "offline benchmark on the official metric."),
    ]},
    4: {85: [
        H("How ranking decisions are explained"),
        B("Each candidate's 1-2 sentence justification is assembled from the "
          "exact factors that set its score: real technologies found in their "
          "career text, real signal values, and honest concerns (relocation, "
          "notice period, gaps)."),
        H("Preventing hallucination"),
        B("Reasoning cites only facts present in the profile. An automated "
          "audit re-checks every number, employer, date and technology on all "
          "100 rows against the source JSON; all six official checks pass."),
        H("Handling suspicious or low-quality profiles"),
        B("A consistency gate discards internally-impossible profiles "
          "(claimed years > career span, expert skills never used, summary "
          "contradicting the experience field). Zero honeypots in our top-100."),
    ]},
    5: {92: [
        H("JD input to ranked output"),
        P("1.  Parse JD text into a structured spec (bands, location, "
          "competencies, disqualifiers)."),
        P("2.  Pre-compute once, offline: MiniLM embeddings + TF-IDF over the "
          "100K career texts."),
        P("3.  Stream candidates -> title gate -> consistency gate -> "
          "evidence floor."),
        P("4.  Score survivors with the multiplicative factor formula "
          "(tuned weights)."),
        P("5.  Take the top 100; generate fact-checked reasoning per "
          "candidate."),
        P("6.  Emit ranked CSV; validate format and audit reasoning "
          "(one command: make judge)."),
        B("Ranking step: ~14s, CPU-only, no network calls."),
    ]},
    6: {None: [   # System Architecture - add a new box (only a title exists)
        P("JD text  ->  JD Parser  ->  structured spec", 11),
        P("Candidates (100K)  ->  [Title gate]  ->  [Consistency gate]  ->  "
          "[Evidence floor]  ->  17,391 survivors", 11),
        P("Precomputed MiniLM embeddings + TF-IDF  feed  ->", 11),
        P("Scoring:  skill_fit x experience x location x availability x "
          "ownership x penalties", 11),
        P("->  Top 100  ->  Reasoning generator  ->  validated + audited CSV", 11),
        ("Eval harness alongside: offline benchmark + minimax weight tuning + "
         "automated reasoning audit.", 10, False, True),
    ]},
    7: {105: [
        H("Ranking quality (offline benchmark, official composite)"),
        B("Composite 0.988 (worst-case 0.987 across six label hypotheses); "
          "NDCG@10 1.000, MAP 1.000, P@10 1.000, NDCG@50 0.960."),
        B("0 honeypots in top-100 (disqualification threshold is >10%). The "
          "organizers' sample ranking scores 0.041 on the same benchmark."),
        B("+/-20% weight perturbation: top-100 stays 97.8% stable."),
        H("Meeting runtime and compute constraints"),
        B("Reproduced in Docker (python:3.11-slim, --network none, -m 16g): "
          "byte-identical CSV, ~50s wall-clock, <1GB RAM."),
        B("Well inside the limits: 5 min / 16 GB / CPU-only / no network."),
        (" All scores are on our own benchmark built on the official metric; "
         "official scores are hidden until submissions close.", 9, False, True),
    ]},
    8: {112: [
        H("Technologies and why"),
        B("Python 3.11 - pipeline, gates, scoring."),
        B("sentence-transformers (all-MiniLM-L6-v2) - CPU embeddings for "
          "semantic match; small and fast, saturates NDCG@10 at ~4x less "
          "compute than larger models."),
        B("scikit-learn TF-IDF - rare-term lexical signal (FAISS, NDCG) that "
          "embeddings miss."),
        B("NumPy - vectorized scoring over 17K survivors in seconds."),
        B("Docker - reproducible, network-isolated evaluation environment."),
        B("Streamlit - hosted sandbox demo."),
        B("Chosen for CPU-only, offline, sub-5-minute operation; no LLM at "
          "rank time, by design and by rule."),
    ]},
    9: {119: [
        H("Submission assets"),
        B("GitHub repo: " + REPO),
        B("Ranked output: submission.csv (top-100, passes the official "
          "validator)."),
        B("Precomputed artifacts: GitHub Release artifacts-v1 (embeddings, "
          "TF-IDF, JD embedding) + precompute.py to rebuild."),
        B("Sandbox demo: " + SANDBOX),
        B("Reproduce: python rank.py --candidates candidates.jsonl --jd "
          "data/job_description.md --out submission.csv   (or: make judge)."),
        B("Docs: README, experiments.md, docs/design_decisions.md, "
          "docs/how_it_works.html."),
    ]},
}


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)

    # Title slide
    s0 = slides[0]
    by_text = {sh.text_frame.paragraphs[0].text.strip(): sh
               for sh in s0.shapes if sh.has_text_frame}
    append_value(by_text["Team Name :"], "Hack_IT_On")
    append_value(by_text["Team Leader Name :"], "Pratham Modi")
    append_value(by_text["Problem Statement :"],
                 "Intelligent Candidate Discovery & Ranking - rank 100,000 "
                 "candidates against a Senior AI Engineer JD and return a "
                 "trustworthy, explainable top-100 shortlist.")

    for idx, shape_map in CONTENT.items():
        slide = slides[idx]
        boxes = {sh.shape_id: sh for sh in slide.shapes if sh.has_text_frame}
        for shape_id, paras in shape_map.items():
            if shape_id is None:
                # add a fresh text box under the title (slide 6)
                from pptx.util import Inches
                box = slide.shapes.add_textbox(Inches(0.41), Inches(1.5),
                                               Inches(9.3), Inches(3.4))
                fill(box, paras)
            else:
                fill(boxes[shape_id], paras)

    prs.save(SRC)
    print("Filled and saved:", SRC)


if __name__ == "__main__":
    main()
